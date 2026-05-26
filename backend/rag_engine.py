import os
import uuid
from io import BytesIO
from typing import List, cast
from pydantic import BaseModel, Field
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from sentence_transformers import CrossEncoder

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import PromptTemplate
from langchain_core.documents import Document

from config import (
    LLM_MODEL, MAX_LENGTH, TEMPERATURE,
    CHUNK_SIZE_CHILD, CHUNK_OVERLAP_CHILD,
    CHUNK_SIZE_PARENT, CHUNK_OVERLAP_PARENT,
    EMBEDDING_MODEL, FAISS_PATH
)


# ==========================================
# DAY 5 BLUEPRINTS: TYPE-SAFE PYDANTIC SCHEMAS
# ==========================================
class SlideModel(BaseModel):
    title: str = Field(description="The clear, strategic business heading for this individual slide.")
    bullets: List[str] = Field(
        description="A sequential list of 3 to 5 highly concise, impactful bullet point takeaways for the slide body.")


class PresentationModel(BaseModel):
    slides: List[SlideModel] = Field(
        description="A cohesive collection of 5 to 7 slides forming the complete presentation deck layout.")


class RagEngine:
    def __init__(self, db_store: str = FAISS_PATH):
        # Initialize dense embedding encoders locally
        self.embeddings = HuggingFaceEmbeddings(
            model_name=EMBEDDING_MODEL,
            model_kwargs={'device': 'cpu'},
            encode_kwargs={'normalize_embeddings': True}
        )

        # Bind the structural Pydantic blueprint to the Gemini API core configuration
        base_llm = ChatGoogleGenerativeAI(
            model=LLM_MODEL,
            temperature=TEMPERATURE,
            max_output_tokens=MAX_LENGTH
        )
        self.llm = base_llm.with_structured_output(PresentationModel)

        # Shared string delimiters to prevent breaking sentences in half
        custom_separators = ["\n\n", "\n", ". ", " ", ""]

        # Day 3 Hierarchical splitters configuration
        self.parent_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE_PARENT,
            chunk_overlap=CHUNK_OVERLAP_PARENT,
            separators=custom_separators
        )
        self.child_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE_CHILD,
            chunk_overlap=CHUNK_OVERLAP_CHILD,
            separators=custom_separators
        )
        self.db_path = db_store
        self.vector_store = self._load_existing_store()

    def _load_existing_store(self):
        """Loads index matrix files locally from local workspace paths if available."""
        if os.path.exists(self.db_path):
            try:
                print(f"Loading existing FAISS index from {self.db_path}...")
                return FAISS.load_local(
                    self.db_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
            except Exception as e:
                print(f"Failed to load existing index: {e}. Starting fresh.")
                return None
        return None

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Parses target files into flat strings using PyPDF2."""
        text = ""
        pdf_reader = PdfReader(file_path)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return text

    def ingest_pdf(self, file_path: str):
        """Builds relational mappings linking dense child vectors to complete parent context strings."""
        raw_text = self.extract_text_from_pdf(file_path)
        if not raw_text.strip():
            raise ValueError("The uploaded PDF seems to be empty or unscannable.")

        parent_texts = self.parent_splitter.split_text(raw_text)
        child_documents = []

        for parent_text in parent_texts:
            parent_id = str(uuid.uuid4())
            sub_chunks = self.child_splitter.split_text(parent_text)

            for chunk in sub_chunks:
                # Store full-context parent text strings directly inside child document metadata fields
                doc = Document(
                    page_content=chunk,
                    metadata={
                        "parent_id": parent_id,
                        "parent_content": parent_text
                    }
                )
                child_documents.append(doc)

        # Commit processed LangChain Document objects safely into FAISS storage indices
        if self.vector_store is None:
            self.vector_store = FAISS.from_documents(child_documents, self.embeddings)
        else:
            self.vector_store.add_documents(child_documents)

        self.vector_store.save_local(self.db_path)
        return len(child_documents)

    @staticmethod
    def _resolve_parent_context(retrieved_child_documents):
        """Extracts unique parent paragraphs out of matched child documents to prevent duplicate data."""
        seen_parent_ids = set()
        parent_blocks = []

        for doc in retrieved_child_documents:
            parent_id = doc.metadata.get("parent_id")
            parent_content = doc.metadata.get("parent_content")

            if parent_id and parent_id not in seen_parent_ids:
                seen_parent_ids.add(parent_id)
                parent_blocks.append(parent_content)

        return "\n\n--- NEXT SECTION BOUNDARY ---\n\n".join(parent_blocks)

    @staticmethod
    def _get_slide_prompt():
        """Returns standard instructions for clean context mapping."""
        template = """
    You are an expert executive business analyst. Your job is to extract high-value insights from the context below and structure them into a professional, cohesive business presentation deck based on the request: "{question}".

    Guidelines:
    - Slide 1 must always be an executive Title Slide.
    - Every slide must focus on a distinct business pillar, timeline step, metric, or strategic insight.
    - Provide exactly 3 to 5 highly concise, impactful bullet points per content slide.

    Context:
    {context}

    Question/Topic Request:
    {question}
    """
        return PromptTemplate(template=template, input_variables=["context", "question"])

    def generate_presentation_data(self, user_question: str) -> PresentationModel:
        """Executes full Two-Stage Retrieval (FAISS + CrossEncoder) and surfaces validated objects."""
        if self.vector_store is None:
            raise ValueError("No documents have been indexed yet. Please upload a PDF file first.")

        # Stage 1: Dynamic Vector candidate retrieval pass
        retriever = self.vector_store.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 15}
        )
        retrieved_child_docs = retriever.invoke(user_question)

        if not retrieved_child_docs:
            raise ValueError("No matching context found in the uploaded documents.")

        # Stage 2: Deep simultaneous query-passage Cross-Encoder evaluation
        reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2", device="cpu")
        pairs = [[user_question, doc.page_content] for doc in retrieved_child_docs]
        scores = reranker.predict(pairs)

        for doc, score in zip(retrieved_child_docs, scores):
            doc.metadata["rerank_score"] = float(score)

        ranked_child_docs = sorted(
            retrieved_child_docs,
            key=lambda x: x.metadata["rerank_score"],
            reverse=True
        )

        # Stage 3: Isolate top candidates and map back to parent content
        top_ranked_docs = ranked_child_docs[:5]
        optimized_parent_context = self._resolve_parent_context(top_ranked_docs)

        # Stage 4: Structure validation output generation pass
        prompt_obj = self._get_slide_prompt()
        formatted_prompt = prompt_obj.format(context=optimized_parent_context, question=user_question)

        # Wrap with explicit cast() to completely satisfy static IDE compiler check tools
        return cast(PresentationModel, self.llm.invoke(formatted_prompt))

    @staticmethod
    def export_to_pptx(presentation_data: PresentationModel) -> BytesIO:
        """Transforms pre-validated Pydantic structures natively into presentation file streams."""
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Standard Widescreen 16:9 Configuration Canvas
        prs.slide_height = Inches(7.5)

        # Safe Object Property Mapping: Removes all brittle text splitting or regex functions
        for slide_data in presentation_data.slides:
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # Draw Title Elements
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_data.title
            p_title.font.size = Pt(38)
            p_title.font.bold = True
            p_title.font.name = "Arial"

            # Draw Content Bullet Elements
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            for i, bullet_text in enumerate(slide_data.bullets):
                p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
                p.text = bullet_text
                p.font.size = Pt(18)
                p.font.name = "Arial"
                p.level = 0
                p.space_after = Pt(14)

        output_stream = BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream